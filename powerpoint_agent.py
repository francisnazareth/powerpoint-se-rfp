"""
PowerPoint AI Agent
A Python AI agent that creates PowerPoint slides using Azure building blocks.
Built with Microsoft Agent Framework and GitHub models.
"""

import asyncio
import os
from typing import Annotated, Dict, Any, List
from datetime import datetime
import json

from agent_framework import ChatAgent
from agent_framework.openai import OpenAIChatClient
from openai import AsyncOpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE


class PowerPointAgent:
    """AI Agent for creating PowerPoint presentations with Azure building blocks."""
    
    def __init__(self, github_token: str, model_id: str = "openai/gpt-4.1"):
        """Initialize the PowerPoint agent.
        
        Args:
            github_token: GitHub personal access token for model access
            model_id: Model ID to use (default: openai/gpt-4.1)
        """
        self.github_token = github_token
        self.model_id = model_id
        self.agent = None
        
        # Azure building blocks definitions
        self.azure_building_blocks = {
            "compute": {
                "virtual_machines": "Azure Virtual Machines (VMs) provide on-demand, scalable computing resources",
                "app_service": "Azure App Service hosts web apps, REST APIs, and mobile backends",
                "azure_functions": "Azure Functions enables serverless compute for event-driven applications",
                "container_instances": "Azure Container Instances run containers without managing servers",
                "kubernetes_service": "Azure Kubernetes Service (AKS) manages containerized applications"
            },
            "storage": {
                "blob_storage": "Azure Blob Storage stores massive amounts of unstructured object data",
                "disk_storage": "Azure Disk Storage provides high-performance, durable block storage",
                "files": "Azure Files offers fully managed file shares in the cloud",
                "queue_storage": "Azure Queue Storage provides messaging between application components",
                "table_storage": "Azure Table Storage stores structured NoSQL data"
            },
            "networking": {
                "virtual_network": "Azure Virtual Network enables secure communication between Azure resources",
                "load_balancer": "Azure Load Balancer distributes incoming traffic across healthy VMs",
                "application_gateway": "Azure Application Gateway provides application-level routing and load balancing",
                "vpn_gateway": "Azure VPN Gateway connects on-premises networks to Azure",
                "express_route": "Azure ExpressRoute creates private connections to Azure datacenters"
            },
            "database": {
                "sql_database": "Azure SQL Database provides managed relational database service",
                "cosmos_db": "Azure Cosmos DB offers globally distributed, multi-model database service",
                "postgresql": "Azure Database for PostgreSQL provides managed PostgreSQL service",
                "mysql": "Azure Database for MySQL offers managed MySQL service",
                "redis_cache": "Azure Cache for Redis provides in-memory data caching"
            },
            "ai_ml": {
                "cognitive_services": "Azure Cognitive Services provides AI capabilities via REST APIs",
                "machine_learning": "Azure Machine Learning enables building and deploying ML models",
                "ai_search": "Azure AI Search provides full-text search capabilities",
                "openai_service": "Azure OpenAI Service offers OpenAI models with enterprise security",
                "bot_service": "Azure Bot Service builds conversational AI experiences"
            },
            "security": {
                "key_vault": "Azure Key Vault securely stores and manages secrets, keys, and certificates",
                "active_directory": "Azure Active Directory provides identity and access management",
                "security_center": "Azure Security Center provides unified security management",
                "sentinel": "Azure Sentinel offers cloud-native SIEM and SOAR capabilities",
                "firewall": "Azure Firewall provides network security filtering"
            },
            "monitoring": {
                "monitor": "Azure Monitor provides comprehensive monitoring for applications and infrastructure",
                "log_analytics": "Azure Log Analytics collects and analyzes log data",
                "application_insights": "Azure Application Insights monitors live applications",
                "service_health": "Azure Service Health provides insights into Azure service issues",
                "advisor": "Azure Advisor provides personalized recommendations"
            }
        }
    
    async def initialize_agent(self):
        """Initialize the AI agent with tools."""
        openai_client = AsyncOpenAI(
            base_url="https://models.github.ai/inference",
            api_key=self.github_token,
        )
        
        chat_client = OpenAIChatClient(
            async_client=openai_client,
            model_id=self.model_id
        )
        
        self.agent = ChatAgent(
            chat_client=chat_client,
            name="PowerPointAgent",
            instructions="""You are a PowerPoint creation expert specializing in Azure architecture.
            
            Your role is to:
            1. Analyze requests for PowerPoint slides about Azure solutions
            2. Suggest appropriate Azure building blocks based on the request
            3. Create structured slide content with clear explanations
            4. Generate professional PowerPoint presentations
            
            When creating slides, focus on:
            - Clear, concise content
            - Logical architecture flow
            - Best practices for Azure services
            - Professional formatting and layout
            
            Use the available tools to create PowerPoint files based on user requests.""",
            tools=[
                self.get_azure_services,
                self.create_powerpoint_slide,
                self.create_architecture_diagram,
                self.save_presentation
            ]
        )
    
    def get_azure_services(
        self, 
        category: Annotated[str, "Category of Azure services (compute, storage, networking, database, ai_ml, security, monitoring)"]
    ) -> str:
        """Get information about Azure services in a specific category."""
        if category.lower() in self.azure_building_blocks:
            services = self.azure_building_blocks[category.lower()]
            result = f"Azure {category.title()} Services:\n\n"
            for service, description in services.items():
                result += f"• {service.replace('_', ' ').title()}: {description}\n"
            return result
        else:
            available_categories = ", ".join(self.azure_building_blocks.keys())
            return f"Category '{category}' not found. Available categories: {available_categories}"
    
    def create_powerpoint_slide(
        self,
        title: Annotated[str, "Title of the slide"],
        content: Annotated[str, "Main content for the slide"],
        slide_type: Annotated[str, "Type of slide: title, content, or architecture"] = "content"
    ) -> str:
        """Create a PowerPoint slide with the specified content."""
        try:
            # Initialize presentation if it doesn't exist
            if not hasattr(self, 'presentation'):
                self.presentation = Presentation()
                # Remove the default slide
                if len(self.presentation.slides) > 0:
                    slide_to_remove = self.presentation.slides[0]
                    rId = self.presentation.slides.slides._element.index(slide_to_remove._element)
                    self.presentation.part.drop_rel(self.presentation.slides._sld_id_lst[rId].rId)
                    del self.presentation.slides._sld_id_lst[rId]
            
            # Add slide based on type
            if slide_type.lower() == "title":
                slide_layout = self.presentation.slide_layouts[0]  # Title slide
                slide = self.presentation.slides.add_slide(slide_layout)
                slide.shapes.title.text = title
                if slide.shapes.placeholders[1]:  # Subtitle
                    slide.shapes.placeholders[1].text = content
            else:
                slide_layout = self.presentation.slide_layouts[1]  # Title and content
                slide = self.presentation.slides.add_slide(slide_layout)
                slide.shapes.title.text = title
                
                # Add content to the body
                content_placeholder = slide.shapes.placeholders[1]
                content_placeholder.text = content
                
                # Format the text
                for paragraph in content_placeholder.text_frame.paragraphs:
                    paragraph.font.size = Pt(14)
                    paragraph.font.name = 'Segoe UI'
            
            slide_number = len(self.presentation.slides)
            return f"Successfully created slide {slide_number}: '{title}'"
            
        except Exception as e:
            return f"Error creating slide: {str(e)}"
    
    def create_architecture_diagram(
        self,
        title: Annotated[str, "Title of the architecture slide"],
        components: Annotated[List[str], "List of Azure components to include in the diagram"],
        description: Annotated[str, "Description of the architecture"]
    ) -> str:
        """Create an architecture diagram slide with Azure components."""
        try:
            # Initialize presentation if it doesn't exist
            if not hasattr(self, 'presentation'):
                self.presentation = Presentation()
                if len(self.presentation.slides) > 0:
                    slide_to_remove = self.presentation.slides[0]
                    rId = self.presentation.slides.slides._element.index(slide_to_remove._element)
                    self.presentation.part.drop_rel(self.presentation.slides._sld_id_lst[rId].rId)
                    del self.presentation.slides._sld_id_lst[rId]
            
            # Create blank slide for architecture diagram
            slide_layout = self.presentation.slide_layouts[6]  # Blank layout
            slide = self.presentation.slides.add_slide(slide_layout)
            
            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = Pt(24)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            
            # Create simple component boxes
            start_x = 1.5
            start_y = 2.0
            box_width = 2.0
            box_height = 1.0
            spacing = 0.5
            
            colors = [
                RGBColor(0, 120, 212),    # Azure blue
                RGBColor(0, 188, 140),    # Green
                RGBColor(255, 140, 0),    # Orange
                RGBColor(232, 17, 35),    # Red
                RGBColor(136, 23, 152),   # Purple
                RGBColor(16, 110, 190),   # Dark blue
            ]
            
            for i, component in enumerate(components[:6]):  # Limit to 6 components
                x_pos = start_x + (i % 3) * (box_width + spacing)
                y_pos = start_y + (i // 3) * (box_height + spacing)
                
                # Create rectangle shape
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(x_pos), Inches(y_pos),
                    Inches(box_width), Inches(box_height)
                )
                
                # Set fill color
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = colors[i % len(colors)]
                
                # Add text
                text_frame = shape.text_frame
                text_frame.text = component.replace('_', ' ').title()
                text_frame.text_frame_format.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                paragraph = text_frame.paragraphs[0]
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.font.bold = True
                paragraph.font.size = Pt(12)
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            
            # Add description box
            desc_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(5.5), Inches(9), Inches(1.5)
            )
            desc_frame = desc_box.text_frame
            desc_frame.text = description
            desc_frame.paragraphs[0].font.size = Pt(14)
            desc_frame.word_wrap = True
            
            slide_number = len(self.presentation.slides)
            return f"Successfully created architecture diagram slide {slide_number}: '{title}' with {len(components)} components"
            
        except Exception as e:
            return f"Error creating architecture diagram: {str(e)}"
    
    def save_presentation(
        self,
        filename: Annotated[str, "Name of the PowerPoint file to save (without .pptx extension)"]
    ) -> str:
        """Save the PowerPoint presentation to a file."""
        try:
            if not hasattr(self, 'presentation'):
                return "No presentation to save. Create slides first."
            
            # Ensure filename ends with .pptx
            if not filename.endswith('.pptx'):
                filename += '.pptx'
            
            # Save to the current directory
            filepath = os.path.join(os.getcwd(), filename)
            self.presentation.save(filepath)
            
            slide_count = len(self.presentation.slides)
            return f"Successfully saved presentation '{filename}' with {slide_count} slides to {filepath}"
            
        except Exception as e:
            return f"Error saving presentation: {str(e)}"
    
    async def create_presentation(self, user_request: str) -> str:
        """Process user request to create a PowerPoint presentation."""
        if not self.agent:
            await self.initialize_agent()
        
        thread = self.agent.get_new_thread()
        
        print("🤖 PowerPoint Agent: ", end="", flush=True)
        response_text = ""
        
        async for chunk in self.agent.run_stream(user_request, thread=thread):
            if chunk.text:
                print(chunk.text, end="", flush=True)
                response_text += chunk.text
        
        print("\n")
        return response_text


# Example usage and main function
async def main():
    """Main function to demonstrate the PowerPoint agent."""
    
    # Get GitHub token from environment variable
    github_token = os.getenv('GITHUB_TOKEN')
    if not github_token:
        print("❌ Please set the GITHUB_TOKEN environment variable with your GitHub Personal Access Token")
        print("   You can create one at: https://github.com/settings/tokens")
        return
    
    # Initialize the agent
    agent = PowerPointAgent(github_token)
    
    print("🚀 PowerPoint AI Agent initialized!")
    print("📋 Available Azure building block categories:")
    for category in agent.azure_building_blocks.keys():
        print(f"   • {category}")
    print()
    
    # Example requests
    sample_requests = [
        "Create a presentation about a web application architecture on Azure with App Service, SQL Database, and Key Vault",
        "Build slides for an AI-powered chatbot solution using Azure OpenAI, Cosmos DB, and Application Gateway",
        "Generate a presentation about Azure data analytics pipeline with Storage Account, Data Factory, and Synapse Analytics"
    ]
    
    print("📝 Sample requests you can try:")
    for i, request in enumerate(sample_requests, 1):
        print(f"   {i}. {request}")
    print()
    
    # Interactive mode
    while True:
        try:
            user_input = input("Enter your PowerPoint request (or 'quit' to exit): ").strip()
            
            if user_input.lower() in ['quit', 'exit', 'q']:
                print("👋 Goodbye!")
                break
            
            if not user_input:
                continue
            
            print(f"\n🔄 Processing your request...\n")
            
            # Create presentation based on user request
            result = await agent.create_presentation(user_input)
            
            print(f"\n✅ Request completed!")
            
        except KeyboardInterrupt:
            print("\n👋 Goodbye!")
            break
        except Exception as e:
            print(f"❌ Error: {str(e)}")


if __name__ == "__main__":
    asyncio.run(main())