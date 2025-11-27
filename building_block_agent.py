"""
Building Block PowerPoint Agent
Creates single-slide PowerPoint presentations following the ConceptualArchitecture-Samples.pptx format.
Built with Microsoft Agent Framework and specific Azure service recommendations.
"""

import asyncio
import os
from typing import Annotated, Dict, Any, List
from datetime import datetime
import json

from agent_framework import ChatAgent
from agent_framework.openai import OpenAIChatClient
from openai import AsyncOpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE


class BuildingBlockAgent:
    """AI Agent for creating building block architecture slides based on requirements."""
    
    def __init__(self, github_token: str, model_id: str = "openai/gpt-4.1"):
        """Initialize the Building Block agent.
        
        Args:
            github_token: GitHub personal access token for model access
            model_id: Model ID to use (default: openai/gpt-4.1)
        """
        self.github_token = github_token
        self.model_id = model_id
        self.agent = None
        
        # Azure service recommendations by category
        self.service_recommendations = {
            "ai_analytics": {
                "primary": ["Azure OpenAI", "Microsoft Fabric", "Azure Databricks"],
                "supporting": ["Azure AI Services", "Azure Synapse Analytics", "Azure ML", "Power BI", "Azure AI Search"]
            },
            "web_application": {
                "primary": ["Azure Web Apps", "Azure Container Apps", "Azure Kubernetes Service"],
                "supporting": ["Azure App Service", "Azure Front Door", "Azure Application Gateway", "Azure Load Balancer"]
            },
            "data_platform": {
                "primary": ["Azure SQL Database", "Azure Cosmos DB", "Azure Storage"],
                "supporting": ["Azure Data Factory", "Azure Data Lake", "Azure Synapse", "Azure Purview"]
            },
            "integration": {
                "primary": ["Azure API Management", "Azure Service Bus", "Azure Logic Apps"],
                "supporting": ["Azure Event Grid", "Azure Event Hub", "Function Apps", "Power Automate"]
            },
            "security": {
                "primary": ["Microsoft Entra ID", "Azure Key Vault", "Microsoft Sentinel"],
                "supporting": ["Azure Firewall", "Microsoft Defender", "Azure Policy", "Azure Monitor"]
            },
            "infrastructure": {
                "primary": ["Azure Virtual Networks", "Azure Virtual Machines", "Azure Backup"],
                "supporting": ["Azure DevOps", "Azure Monitor", "Azure Policy and Compliance"]
            }
        }
        
        # Building block colors (based on Azure color palette)
        self.building_block_colors = {
            "ai_analytics": RGBColor(138, 43, 226),      # Purple for AI/ML
            "web_application": RGBColor(0, 120, 212),    # Azure Blue for Web
            "data_platform": RGBColor(0, 188, 140),      # Teal for Data
            "integration": RGBColor(255, 140, 0),        # Orange for Integration
            "security": RGBColor(232, 17, 35),           # Red for Security
            "infrastructure": RGBColor(16, 110, 190),    # Dark Blue for Infrastructure
            "default": RGBColor(68, 68, 68)              # Gray for others
        }
    
    async def initialize_agent(self):
        """Initialize the AI agent with tools."""
        openai_client = AsyncOpenAI(
            base_url="https://models.github.ai/inference",
            api_key=self.github_token,
        )
        
        chat_client = OpenAIChatClient(
            async_client=openai_client,
            model_id=self.model_id
        )
        
        self.agent = ChatAgent(
            chat_client=chat_client,
            name="BuildingBlockAgent",
            instructions="""You are an expert Azure solution architect specializing in creating building block diagrams.

Your role is to:
1. Analyze user requirements and identify the main solution categories needed
2. Recommend appropriate Azure services based on predefined service recommendations
3. Create a single comprehensive building block slide showing the complete architecture
4. Organize services into logical building blocks with clear relationships

For different requirement types, use these specific service recommendations:

AI and Analytics requirements:
- Primary: Azure OpenAI, Microsoft Fabric, Azure Databricks
- Supporting: Azure AI Services, Azure Synapse Analytics, Azure ML, Power BI

Web Application requirements:
- Primary: Azure Web Apps, Azure Container Apps, Azure Kubernetes Service
- Supporting: Azure App Service, Azure Front Door, Azure Application Gateway

Always create a single slide with all necessary building blocks clearly organized and labeled.
Use the building block format with colored boxes representing different service categories.""",
            tools=[
                self.analyze_requirements,
                self.get_service_recommendations,
                self.create_building_block_slide,
                self.save_presentation
            ]
        )
    
    def analyze_requirements(
        self,
        requirements: Annotated[str, "User requirements for the solution"]
    ) -> str:
        """Analyze requirements and identify needed Azure service categories."""
        requirements_lower = requirements.lower()
        needed_categories = []
        
        # Check for specific layer mentions (architectural layers)
        if "user experience" in requirements_lower or "ux" in requirements_lower or "frontend" in requirements_lower:
            needed_categories.append("web_application")
        
        if "application layer" in requirements_lower or "app layer" in requirements_lower:
            needed_categories.append("web_application")
        
        if "data and intelligence" in requirements_lower or "data intelligence" in requirements_lower or "analytics" in requirements_lower:
            needed_categories.append("ai_analytics")
            needed_categories.append("data_platform")
        
        if "integration layer" in requirements_lower:
            needed_categories.append("integration")
        
        # Check for AI/Analytics requirements
        ai_keywords = ["ai", "analytics", "machine learning", "ml", "data science", "openai", "chatbot", "intelligent", "prediction"]
        if any(keyword in requirements_lower for keyword in ai_keywords):
            needed_categories.append("ai_analytics")
        
        # Check for Web Application requirements
        web_keywords = ["web", "website", "portal", "api", "frontend", "backend", "application", "app"]
        if any(keyword in requirements_lower for keyword in web_keywords):
            needed_categories.append("web_application")
        
        # Check for Data Platform requirements
        data_keywords = ["database", "data", "storage", "sql", "cosmos", "warehouse", "lake"]
        if any(keyword in requirements_lower for keyword in data_keywords):
            needed_categories.append("data_platform")
        
        # Check for Integration requirements
        integration_keywords = ["integration", "api", "messaging", "event", "workflow", "automation"]
        if any(keyword in requirements_lower for keyword in integration_keywords):
            needed_categories.append("integration")
        
        # Check for Security requirements
        security_keywords = ["security", "authentication", "authorization", "identity", "firewall", "compliance"]
        if any(keyword in requirements_lower for keyword in security_keywords):
            needed_categories.append("security")
        
        # Always include infrastructure for comprehensive solutions
        needed_categories.append("infrastructure")
        
        # Remove duplicates while preserving order
        needed_categories = list(dict.fromkeys(needed_categories))
        
        result = f"Analysis of requirements: '{requirements}'\n\nIdentified solution categories:\n"
        for category in needed_categories:
            result += f"• {category.replace('_', ' ').title()}\n"
        
        return result
    
    def get_service_recommendations(
        self,
        categories: Annotated[List[str], "List of categories to get service recommendations for"]
    ) -> str:
        """Get Azure service recommendations for specified categories."""
        result = "Azure Service Recommendations:\n\n"
        
        for category in categories:
            if category in self.service_recommendations:
                services = self.service_recommendations[category]
                result += f"{category.replace('_', ' ').title()}:\n"
                result += f"  Primary services: {', '.join(services['primary'])}\n"
                result += f"  Supporting services: {', '.join(services['supporting'])}\n\n"
        
        return result
    
    def create_building_block_slide(
        self,
        title: Annotated[str, "Title for the building block slide"],
        requirements: Annotated[str, "Original user requirements"],
        categories: Annotated[List[str], "List of solution categories to include"]
    ) -> str:
        """Create a building block architecture slide."""
        try:
            # Initialize presentation
            if not hasattr(self, 'presentation'):
                self.presentation = Presentation()
                # Remove default slide
                if len(self.presentation.slides) > 0:
                    slide_to_remove = self.presentation.slides[0]
                    rId = self.presentation.slides.slides._element.index(slide_to_remove._element)
                    self.presentation.part.drop_rel(self.presentation.slides._sld_id_lst[rId].rId)
                    del self.presentation.slides._sld_id_lst[rId]
            
            # Create blank slide for building blocks
            slide_layout = self.presentation.slide_layouts[6]  # Blank layout
            slide = self.presentation.slides.add_slide(slide_layout)
            
            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = Pt(20)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            
            # Create building blocks layout - dynamic sizing based on number of categories
            num_categories = len([cat for cat in categories if cat in self.service_recommendations])
            if num_categories <= 3:
                blocks_per_row = num_categories
                block_width = 2.8
                block_height = 2.0
            elif num_categories <= 6:
                blocks_per_row = 3
                block_width = 2.5
                block_height = 1.8
            else:
                blocks_per_row = 3
                block_width = 2.2
                block_height = 1.5
            
            start_x = 0.5
            start_y = 1.8
            spacing_x = 0.3
            spacing_y = 0.4
            
            row = 0
            col = 0
            
            for i, category in enumerate(categories):
                if category in self.service_recommendations:
                    # Calculate position
                    x_pos = start_x + col * (block_width + spacing_x)
                    y_pos = start_y + row * (block_height + spacing_y)
                    
                    # Create building block rectangle
                    block_shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(x_pos), Inches(y_pos),
                        Inches(block_width), Inches(block_height)
                    )
                    
                    # Set color based on category
                    fill = block_shape.fill
                    fill.solid()
                    if category in self.building_block_colors:
                        fill.fore_color.rgb = self.building_block_colors[category]
                    else:
                        fill.fore_color.rgb = self.building_block_colors["default"]
                    
                    # Add category title and services
                    text_frame = block_shape.text_frame
                    text_frame.clear()
                    text_frame.margin_top = Inches(0.1)
                    text_frame.margin_bottom = Inches(0.1)
                    text_frame.margin_left = Inches(0.1)
                    text_frame.margin_right = Inches(0.1)
                    text_frame.word_wrap = True
                    
                    # Category header
                    p_title = text_frame.paragraphs[0]
                    p_title.text = category.replace('_', ' ').title()
                    p_title.font.bold = True
                    p_title.font.size = Pt(11)
                    p_title.font.color.rgb = RGBColor(255, 255, 255)
                    p_title.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                    
                    # Add primary services
                    services = self.service_recommendations[category]
                    for service in services['primary'][:4]:  # Show top 4 primary services
                        p_service = text_frame.add_paragraph()
                        p_service.text = f"• {service}"
                        p_service.font.size = Pt(8)
                        p_service.font.color.rgb = RGBColor(255, 255, 255)
                        p_service.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                        p_service.space_after = Pt(2)
                    
                    # Move to next position
                    col += 1
                    if col >= blocks_per_row:
                        col = 0
                        row += 1
            
            # Add cross-cutting security and infrastructure services as individual black boxes
            security_y_pos = start_y + ((num_categories - 1) // blocks_per_row + 1) * (block_height + spacing_y) + 0.3
            
            # Cross-cutting services
            cross_cutting_services = [
                "Azure Policy and Compliance", "Azure Firewall and DDoS", "Microsoft Sentinel and Defender",
                "Encryption", "Azure Monitor", "Azure Backup and BCDR", "Microsoft Entra ID"
            ]
            
            # Create individual rectangular boxes for cross-cutting services
            box_width = 1.25
            box_height = 0.6
            start_x_cross = 0.5
            spacing_x_cross = 0.05
            
            for i, service in enumerate(cross_cutting_services):
                x_pos = start_x_cross + i * (box_width + spacing_x_cross)
                
                # Create rectangular box (not rounded)
                service_box = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x_pos), Inches(security_y_pos),
                    Inches(box_width), Inches(box_height)
                )
                
                # Set black background
                fill = service_box.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0, 0, 0)  # Pure black background
                
                # Add service text to the box
                service_text_frame = service_box.text_frame
                service_text_frame.clear()
                service_text_frame.margin_top = Inches(0.05)
                service_text_frame.margin_bottom = Inches(0.05)
                service_text_frame.margin_left = Inches(0.05)
                service_text_frame.margin_right = Inches(0.05)
                service_text_frame.word_wrap = True
                
                # Service text
                p_service = service_text_frame.paragraphs[0]
                p_service.text = service
                p_service.font.size = Pt(8)
                p_service.font.color.rgb = RGBColor(255, 255, 255)  # White text on black background
                p_service.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            
            # Add header label above the cross-cutting services
            header_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(security_y_pos - 0.3), Inches(9), Inches(0.25)
            )
            header_frame = header_box.text_frame
            header_frame.text = "Cross-cutting Security and Infrastructure Services"
            header_frame.paragraphs[0].font.bold = True
            header_frame.paragraphs[0].font.size = Pt(11)
            header_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            
            # Add requirements summary at bottom
            req_y_pos = security_y_pos + box_height + 0.2
            req_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(req_y_pos), Inches(9), Inches(0.5)
            )
            req_frame = req_box.text_frame
            req_frame.text = f"Requirements: {requirements}"
            req_frame.paragraphs[0].font.size = Pt(10)
            req_frame.paragraphs[0].font.italic = True
            req_frame.word_wrap = True
            
            slide_number = len(self.presentation.slides)
            return f"Successfully created building block slide: '{title}' with {len(categories)} building blocks"
            
        except Exception as e:
            return f"Error creating building block slide: {str(e)}"
    
    def save_presentation(
        self,
        filename: Annotated[str, "Name of the PowerPoint file to save (without .pptx extension)"]
    ) -> str:
        """Save the PowerPoint presentation to a file."""
        try:
            if not hasattr(self, 'presentation'):
                return "No presentation to save. Create slides first."
            
            # Ensure filename ends with .pptx
            if not filename.endswith('.pptx'):
                filename += '.pptx'
            
            # Save to the current directory
            filepath = os.path.join(os.getcwd(), filename)
            self.presentation.save(filepath)
            
            slide_count = len(self.presentation.slides)
            return f"Successfully saved presentation '{filename}' with {slide_count} slides to {filepath}"
            
        except Exception as e:
            return f"Error saving presentation: {str(e)}"
    
    async def create_building_block_presentation(self, requirements: str) -> str:
        """Process user requirements to create a building block presentation."""
        if not self.agent:
            await self.initialize_agent()
        
        thread = self.agent.get_new_thread()
        
        prompt = f"""Create a SINGLE building block architecture slide for these requirements:

Requirements: {requirements}

IMPORTANT: Create only ONE slide with ALL building blocks on it. Do not create separate slides for each layer or category.

Steps to follow:
1. Call analyze_requirements to identify ALL needed solution categories from the requirements
2. Call get_service_recommendations for those categories  
3. Call create_building_block_slide ONCE with ALL categories to create a single comprehensive slide
4. Call save_presentation to save the file

The slide should show all building blocks (layers/categories) organized on one slide with their respective Azure services.
Title the slide "Solution Architecture Building Blocks" and include ALL identified categories in a single slide layout."""
        
        print("🏗️ Building Block Agent: ", end="", flush=True)
        response_text = ""
        
        async for chunk in self.agent.run_stream(prompt, thread=thread):
            if chunk.text:
                print(chunk.text, end="", flush=True)
                response_text += chunk.text
        
        print("\n")
        return response_text


# Main interactive function
async def main():
    """Main function for the Building Block agent."""
    
    # Get GitHub token from environment variable
    github_token = os.getenv('GITHUB_TOKEN')
    if not github_token:
        print("❌ Please set the GITHUB_TOKEN environment variable")
        return
    
    # Initialize the agent
    agent = BuildingBlockAgent(github_token, model_id="openai/gpt-4o-mini")  # Using faster model
    
    print("🏗️ Building Block PowerPoint Agent")
    print("=" * 50)
    print("Creates single-slide architecture presentations following your building block format")
    print()
    
    # Show service recommendations
    print("📋 Service Recommendations by Category:")
    print("AI & Analytics: Azure OpenAI, Microsoft Fabric, Azure Databricks")
    print("Web Applications: Azure Web Apps, Azure Container Apps, Azure Kubernetes Service")
    print("Data Platform: Azure SQL Database, Azure Cosmos DB, Azure Storage")
    print("Integration: Azure API Management, Azure Service Bus, Azure Logic Apps")
    print("Security: Microsoft Entra ID, Azure Key Vault, Microsoft Sentinel")
    print()
    
    # Example requirements
    examples = [
        "AI-powered customer service solution with chat capabilities and analytics",
        "Modern web application with microservices architecture and database",
        "Data analytics platform for business intelligence and reporting",
        "E-commerce platform with payment processing and inventory management"
    ]
    
    print("💡 Example requirements:")
    for i, example in enumerate(examples, 1):
        print(f"   {i}. {example}")
    print()
    
    # Interactive loop
    while True:
        try:
            requirements = input("Enter your solution requirements (or 'quit' to exit): ").strip()
            
            if requirements.lower() in ['quit', 'exit', 'q']:
                print("👋 Goodbye!")
                break
            
            if not requirements:
                continue
            
            print(f"\n🔄 Creating building block slide for: {requirements}\n")
            
            # Create presentation
            result = await agent.create_building_block_presentation(requirements)
            
            print(f"\n✅ Building block slide created!")
            print("Check your current directory for the generated .pptx file.")
            print()
            
        except KeyboardInterrupt:
            print("\n👋 Goodbye!")
            break
        except Exception as e:
            print(f"❌ Error: {str(e)}")


if __name__ == "__main__":
    asyncio.run(main())