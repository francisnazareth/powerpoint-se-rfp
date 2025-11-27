"""
Direct Building Block Generator
Creates building block slides directly without AI agent intermediate steps
"""

import os
from typing import List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE


class DirectBuildingBlockGenerator:
    """Direct generator for building block slides."""
    
    def __init__(self):
        # Conceptual building blocks (no Azure product names)
        self.conceptual_blocks = {
            "web_application": {
                "name": "User Experience & Application Layer",
                "concepts": ["Web Portals", "Mobile Applications", "User Interfaces", "Frontend Services"]
            },
            "ai_analytics": {
                "name": "Data & Intelligence Layer", 
                "concepts": ["AI/ML Models", "Analytics Engine", "Data Processing", "Intelligent Services"]
            },
            "data_platform": {
                "name": "Data Platform Layer",
                "concepts": ["Databases", "Data Storage", "Data Lakes", "Data Pipelines"]
            },
            "integration": {
                "name": "Integration & API Layer",
                "concepts": ["API Gateway", "Message Queues", "Workflow Engine", "Event Processing"]
            }
        }
        
        # Azure service recommendations by category
        self.service_recommendations = {
            "web_application": {
                "name": "User Experience & Application Layer",
                "services": ["Azure Web Apps", "Azure Container Apps", "Azure Kubernetes Service", "Azure Front Door"]
            },
            "ai_analytics": {
                "name": "Data & Intelligence Layer", 
                "services": ["Azure OpenAI", "Microsoft Fabric", "Azure Databricks", "Azure AI Services"]
            },
            "data_platform": {
                "name": "Data Platform Layer",
                "services": ["Azure SQL Database", "Azure Cosmos DB", "Azure Storage", "Azure Data Factory"]
            },
            "integration": {
                "name": "Integration & API Layer",
                "services": ["Azure API Management", "Azure Service Bus", "Azure Logic Apps", "Azure Event Grid"]
            }
        }
        
        # Cross-cutting services (always included)
        self.conceptual_cross_cutting = [
            "Security & Compliance", "Network Security", "Monitoring & Defense",
            "Data Encryption", "System Monitoring", "Backup & Recovery", "Identity Management"
        ]
        
        self.azure_cross_cutting_services = [
            "Azure Policy and Compliance", "Azure Firewall and DDoS", "Microsoft Sentinel and Defender",
            "Encryption", "Azure Monitor", "Azure Backup and BCDR", "Microsoft Entra ID"
        ]
        
        # Azure service icons mapping (using Unicode symbols as placeholders for actual icons)
        self.azure_icons = {
            # Web Services
            "Azure Web Apps": "🌐",
            "Azure Container Apps": "📦",
            "Azure Kubernetes Service": "⚙️",
            "Azure Front Door": "🚪",
            "Azure App Service": "🌐",
            "Azure Application Gateway": "🚪",
            "Azure Load Balancer": "⚖️",
            
            # AI/ML Services
            "Azure OpenAI": "🤖",
            "Microsoft Fabric": "🧩",
            "Azure Databricks": "📊",
            "Azure AI Services": "🧠",
            "Azure Synapse Analytics": "📈",
            "Azure ML": "🔬",
            "Power BI": "📊",
            "Azure AI Search": "🔍",
            
            # Data Services
            "Azure SQL Database": "🗄️",
            "Azure Cosmos DB": "🌌",
            "Azure Storage": "💾",
            "Azure Data Factory": "🏭",
            "Azure Data Lake": "🏞️",
            "Azure Synapse": "📈",
            "Azure Purview": "🔍",
            
            # Integration Services
            "Azure API Management": "🔌",
            "Azure Service Bus": "🚌",
            "Azure Logic Apps": "⚡",
            "Azure Event Grid": "📋",
            "Azure Event Hub": "📡",
            "Function Apps": "⚡",
            "Power Automate": "🔄",
            
            # Security Services
            "Microsoft Entra ID": "🔐",
            "Azure Key Vault": "🔑",
            "Microsoft Sentinel": "🛡️",
            "Azure Firewall": "🔥",
            "Microsoft Defender": "🛡️",
            "Azure Policy": "📋",
            "Encryption": "🔒",
            
            # Infrastructure Services
            "Azure Virtual Networks": "🌐",
            "Azure Virtual Machines": "🖥️",
            "Azure Monitor": "📊",
            "Azure DevOps": "🔧",
            "Azure Backup": "💾",
            "Azure Policy and Compliance": "📋",
            "Azure Firewall and DDoS": "🔥",
            "Microsoft Sentinel and Defender": "🛡️",
            "Azure Backup and BCDR": "💾"
        }
        
        # Building block colors
        self.colors = {
            "web_application": RGBColor(0, 120, 212),    # Azure Blue
            "ai_analytics": RGBColor(138, 43, 226),      # Purple  
            "data_platform": RGBColor(0, 188, 140),      # Teal
            "integration": RGBColor(255, 140, 0),        # Orange
            "security": RGBColor(232, 17, 35),           # Red
            "infrastructure": RGBColor(16, 110, 190)     # Dark Blue
        }
    
    def analyze_requirements(self, requirements: str) -> List[str]:
        """Analyze requirements and return needed categories."""
        req_lower = requirements.lower()
        categories = []
        
        # Check for layer mentions
        if any(term in req_lower for term in ["user experience", "ux", "frontend", "application layer", "app layer", "web", "portal"]):
            categories.append("web_application")
        
        if any(term in req_lower for term in ["data and intelligence", "data intelligence", "analytics", "ai", "machine learning"]):
            categories.append("ai_analytics")
            if "data" in req_lower:
                categories.append("data_platform")
        
        if any(term in req_lower for term in ["integration layer", "integration", "api"]):
            categories.append("integration")
        
        # Remove duplicates while preserving order
        # Note: Cross-cutting security and infrastructure services are always added as a separate bar
        return list(dict.fromkeys(categories))
    
    def get_azure_icon_symbol(self, service_name: str) -> str:
        """Get the icon symbol for an Azure service."""
        return self.azure_icons.get(service_name, "🔧")  # Default to gear icon
    
    def add_azure_icon_shape(self, slide, x_pos, y_pos, service_name):
        """Add a visual Azure icon shape (placeholder for actual Azure icons)."""
        # This is a placeholder method that could be enhanced to load actual Azure SVG icons
        # For now, it creates a small colored circle as an icon placeholder
        
        # Get service category to determine color
        service_category = None
        for category, data in self.service_recommendations.items():
            if service_name in data["services"]:
                service_category = category
                break
        
        if service_category and service_category in self.colors:
            icon_color = self.colors[service_category]
        else:
            icon_color = RGBColor(100, 100, 100)  # Default gray
        
        # Create small circle as icon placeholder
        icon_size = 0.15
        icon_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x_pos), Inches(y_pos),
            Inches(icon_size), Inches(icon_size)
        )
        
        # Set icon color
        fill = icon_shape.fill
        fill.solid()
        fill.fore_color.rgb = icon_color
        
        return icon_shape
    
    def load_azure_icon_image(self, slide, x_pos, y_pos, service_name, icon_folder="azure_icons"):
        """Load actual Azure service icon if available (SVG or PNG)."""
        import os
        
        # Check if icon folder exists
        if not os.path.exists(icon_folder):
            print(f"Icon folder '{icon_folder}' not found. Using fallback icons.")
            return None
        
        # Common Azure service icon filename patterns
        possible_filenames = [
            f"{service_name.lower().replace(' ', '-')}.svg",
            f"{service_name.lower().replace(' ', '_')}.svg",
            f"{service_name.lower().replace(' ', '-')}.png",
            f"{service_name.lower().replace(' ', '_')}.png",
            f"{service_name.lower()}.svg",
            f"{service_name.lower()}.png"
        ]
        
        for filename in possible_filenames:
            icon_path = os.path.join(icon_folder, filename)
            if os.path.exists(icon_path):
                try:
                    # Add image to slide
                    icon_shape = slide.shapes.add_picture(
                        icon_path,
                        Inches(x_pos), Inches(y_pos),
                        Inches(0.3), Inches(0.3)  # Small icon size
                    )
                    print(f"Loaded Azure icon: {filename}")
                    return icon_shape
                except Exception as e:
                    print(f"Error loading icon {filename}: {e}")
                    continue
        
        print(f"No icon file found for {service_name}. Using fallback.")
        return None
    
    def create_conceptual_slide(self, prs, categories, requirements):
        """Create the first slide with conceptual building blocks (no Azure product names)."""
        # Create blank slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Conceptual Architecture Building Blocks"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        
        # Calculate layout for conceptual blocks
        num_categories = len(categories)
        if num_categories <= 3:
            blocks_per_row = num_categories
            block_width = 2.8
            block_height = 2.2
        else:
            blocks_per_row = 3
            block_width = 2.5
            block_height = 2.0
        
        start_x = 0.5
        start_y = 1.3
        spacing_x = 0.3
        spacing_y = 0.4
        
        # Create conceptual building blocks
        for i, category in enumerate(categories):
            if category in self.conceptual_blocks:
                row = i // blocks_per_row
                col = i % blocks_per_row
                
                x_pos = start_x + col * (block_width + spacing_x)
                y_pos = start_y + row * (block_height + spacing_y)
                
                # Create block
                block_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(x_pos), Inches(y_pos),
                    Inches(block_width), Inches(block_height)
                )
                
                # Set color (lighter/pastel versions)
                fill = block_shape.fill
                fill.solid()
                fill.fore_color.rgb = self.colors[category]
                
                # Add text
                text_frame = block_shape.text_frame
                text_frame.clear()
                text_frame.margin_top = Inches(0.1)
                text_frame.margin_bottom = Inches(0.1)
                text_frame.margin_left = Inches(0.1)
                text_frame.margin_right = Inches(0.1)
                text_frame.word_wrap = True
                
                # Category title
                p_title = text_frame.paragraphs[0]
                p_title.text = self.conceptual_blocks[category]["name"]
                p_title.font.bold = True
                p_title.font.size = Pt(12)
                p_title.font.color.rgb = RGBColor(255, 255, 255)
                p_title.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                
                # Add conceptual services
                concepts = self.conceptual_blocks[category]["concepts"]
                for concept in concepts[:4]:  # Show top 4 concepts
                    p_concept = text_frame.add_paragraph()
                    p_concept.text = f"• {concept}"
                    p_concept.font.size = Pt(9)
                    p_concept.font.color.rgb = RGBColor(255, 255, 255)
                    p_concept.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    p_concept.space_after = Pt(3)
        
        # Add conceptual cross-cutting services as black boxes
        security_y_pos = start_y + ((num_categories - 1) // blocks_per_row + 1) * (block_height + spacing_y) + 0.3
        
        box_width = 1.25
        box_height = 0.6
        start_x_cross = 0.5
        spacing_x_cross = 0.05
        
        for i, service in enumerate(self.conceptual_cross_cutting):
            x_pos = start_x_cross + i * (box_width + spacing_x_cross)
            
            # Create rectangular box (not rounded)
            service_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x_pos), Inches(security_y_pos),
                Inches(box_width), Inches(box_height)
            )
            
            # Set black background
            fill = service_box.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)  # Pure black background
            
            # Add service text to the box
            service_text_frame = service_box.text_frame
            service_text_frame.clear()
            service_text_frame.margin_top = Inches(0.05)
            service_text_frame.margin_bottom = Inches(0.05)
            service_text_frame.margin_left = Inches(0.05)
            service_text_frame.margin_right = Inches(0.05)
            service_text_frame.word_wrap = True
            
            # Service text
            p_service = service_text_frame.paragraphs[0]
            p_service.text = service
            p_service.font.size = Pt(8)
            p_service.font.color.rgb = RGBColor(255, 255, 255)  # White text on black background
            p_service.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        
        # Add header label above the cross-cutting services
        header_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(security_y_pos - 0.3), Inches(9), Inches(0.25)
        )
        header_frame = header_box.text_frame
        header_frame.text = "Cross-cutting Security and Infrastructure Services"
        header_frame.paragraphs[0].font.bold = True
        header_frame.paragraphs[0].font.size = Pt(11)
        header_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        
        return security_y_pos, box_height

    def create_azure_specific_slide(self, prs, categories, requirements):
        """Create the second slide with Azure-specific services."""
        # Create blank slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Azure Solution Architecture Building Blocks"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        
        # Calculate layout
        num_categories = len(categories)
        if num_categories <= 3:
            blocks_per_row = num_categories
            block_width = 2.8
            block_height = 2.2
        else:
            blocks_per_row = 3
            block_width = 2.5
            block_height = 2.0
        
        start_x = 0.5
        start_y = 1.3
        spacing_x = 0.3
        spacing_y = 0.4
        
        # Create building blocks
        for i, category in enumerate(categories):
            if category in self.service_recommendations:
                row = i // blocks_per_row
                col = i % blocks_per_row
                
                x_pos = start_x + col * (block_width + spacing_x)
                y_pos = start_y + row * (block_height + spacing_y)
                
                # Create block
                block_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(x_pos), Inches(y_pos),
                    Inches(block_width), Inches(block_height)
                )
                
                # Set color
                fill = block_shape.fill
                fill.solid()
                fill.fore_color.rgb = self.colors[category]
                
                # Add text
                text_frame = block_shape.text_frame
                text_frame.clear()
                text_frame.margin_top = Inches(0.1)
                text_frame.margin_bottom = Inches(0.1)
                text_frame.margin_left = Inches(0.1)
                text_frame.margin_right = Inches(0.1)
                text_frame.word_wrap = True
                
                # Category title
                p_title = text_frame.paragraphs[0]
                p_title.text = self.service_recommendations[category]["name"]
                p_title.font.bold = True
                p_title.font.size = Pt(12)
                p_title.font.color.rgb = RGBColor(255, 255, 255)
                p_title.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                
                # Add Azure services with icons
                services = self.service_recommendations[category]["services"]
                for j, service in enumerate(services[:4]):  # Show top 4 services
                    p_service = text_frame.add_paragraph()
                    
                    # Try to load actual Azure icon first, fallback to Unicode symbol
                    icon_loaded = self.load_azure_icon_image(
                        slide, (shape.left / 914400) - 0.35, (shape.top / 914400) + (len(text_frame.paragraphs) * 0.15), service
                    )
                    
                    if icon_loaded:
                        # If we loaded an actual icon, just use the service name
                        p_service.text = service
                    else:
                        # Fallback to Unicode symbol + service name
                        icon_symbol = self.get_azure_icon_symbol(service)
                        p_service.text = f"{icon_symbol} {service}"
                    p_service.font.size = Pt(9)
                    p_service.font.color.rgb = RGBColor(255, 255, 255)
                    p_service.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    p_service.space_after = Pt(3)
        
        # Add Azure cross-cutting services as black boxes
        security_y_pos = start_y + ((num_categories - 1) // blocks_per_row + 1) * (block_height + spacing_y) + 0.3
        
        box_width = 1.25
        box_height = 0.6
        start_x_cross = 0.5
        spacing_x_cross = 0.05
        
        for i, service in enumerate(self.azure_cross_cutting_services):
            x_pos = start_x_cross + i * (box_width + spacing_x_cross)
            
            # Create rectangular box (not rounded)
            service_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x_pos), Inches(security_y_pos),
                Inches(box_width), Inches(box_height)
            )
            
            # Set black background
            fill = service_box.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)  # Pure black background
            
            # Add service text to the box
            service_text_frame = service_box.text_frame
            service_text_frame.clear()
            service_text_frame.margin_top = Inches(0.05)
            service_text_frame.margin_bottom = Inches(0.05)
            service_text_frame.margin_left = Inches(0.05)
            service_text_frame.margin_right = Inches(0.05)
            service_text_frame.word_wrap = True
            
            # Service text with icon
            p_service = service_text_frame.paragraphs[0]
            
            # Try to load actual Azure icon first, fallback to Unicode symbol
            icon_loaded = self.load_azure_icon_image(
                slide, (service_shape.left / 914400), (service_shape.top / 914400) - 0.2, service
            )
            
            if icon_loaded:
                # If we loaded an actual icon, just use the service name
                p_service.text = service
            else:
                # Fallback to Unicode symbol + service name
                icon_symbol = self.get_azure_icon_symbol(service)
                p_service.text = f"{icon_symbol}\n{service}"
            p_service.font.size = Pt(7)
            p_service.font.color.rgb = RGBColor(255, 255, 255)  # White text on black background
            p_service.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        
        # Add header label above the cross-cutting services
        header_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(security_y_pos - 0.3), Inches(9), Inches(0.25)
        )
        header_frame = header_box.text_frame
        header_frame.text = "Cross-cutting Security and Infrastructure Services"
        header_frame.paragraphs[0].font.bold = True
        header_frame.paragraphs[0].font.size = Pt(11)
        header_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        
        return security_y_pos, box_height

    def create_building_block_slide(self, requirements: str, filename: str = None) -> str:
        """Create two building block slides: conceptual and Azure-specific."""
        try:
            # Analyze requirements
            categories = self.analyze_requirements(requirements)
            
            print(f"Identified categories: {[cat.replace('_', ' ').title() for cat in categories]}")
            
            # Create presentation
            prs = Presentation()
            # Remove default slide
            if len(prs.slides) > 0:
                slide_to_remove = prs.slides[0]
                rId = prs.slides.slides._element.index(slide_to_remove._element)
                prs.part.drop_rel(prs.slides._sld_id_lst[rId].rId)
                del prs.slides._sld_id_lst[rId]
            
            # Create first slide: Conceptual building blocks
            print("Creating conceptual building blocks slide...")
            security_y_pos1, box_height1 = self.create_conceptual_slide(prs, categories, requirements)
            
            # Create second slide: Azure-specific building blocks  
            print("Creating Azure-specific building blocks slide...")
            security_y_pos2, box_height2 = self.create_azure_specific_slide(prs, categories, requirements)
            
            # Add requirements at bottom of both slides
            for slide_num, slide in enumerate(prs.slides, 1):
                req_y_pos = security_y_pos2 + box_height2 + 0.2
                req_box = slide.shapes.add_textbox(Inches(0.5), Inches(req_y_pos), Inches(9), Inches(0.5))
                req_frame = req_box.text_frame
                req_frame.text = f"Requirements: {requirements}"
                req_frame.paragraphs[0].font.size = Pt(10)
                req_frame.paragraphs[0].font.italic = True
                req_frame.word_wrap = True
            
            # Save presentation
            if not filename:
                filename = "Building_Blocks_Architecture.pptx"
            elif not filename.endswith('.pptx'):
                filename += '.pptx'
            
            filepath = os.path.join(os.getcwd(), filename)
            prs.save(filepath)
            
            slide_count = len(prs.slides)
            return f"✅ Successfully created {slide_count} slides with {len(categories)} building blocks\n   Slide 1: Conceptual Architecture Building Blocks\n   Slide 2: Azure-Specific Architecture Building Blocks\nSaved as: {filename}"
            
        except Exception as e:
            return f"❌ Error creating slide: {str(e)}"


def main():
    """Main function for direct building block generation."""
    generator = DirectBuildingBlockGenerator()
    
    print("🏗️ Direct Building Block Generator")
    print("=" * 50)
    
    # Test with your example
    test_requirements = "1. User experience layer, 2. Application Layer, 3. Data and intelligence layer, 4. Integration Layer"
    
    print(f"Test Requirements: {test_requirements}")
    print()
    
    result = generator.create_building_block_slide(test_requirements, "Multi_Layer_Architecture")
    print(result)
    
    print("\n" + "=" * 50)
    print("Enter your own requirements:")
    
    while True:
        requirements = input("\nRequirements (or 'quit'): ").strip()
        
        if requirements.lower() in ['quit', 'exit', 'q']:
            break
        
        if requirements:
            result = generator.create_building_block_slide(requirements)
            print(result)


if __name__ == "__main__":
    main()