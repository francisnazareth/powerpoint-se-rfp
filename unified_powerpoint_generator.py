"""
Unified PowerPoint Generator
Combines all PowerPoint generation logic into a single comprehensive file.
Supports both AI-powered and direct generation methods with Azure building blocks.
"""

import asyncio
import os
import sys
from typing import Annotated, Dict, Any, List
from datetime import datetime
import json

# AI Agent imports (optional - will gracefully degrade if not available)
try:
    from agent_framework import ChatAgent
    from agent_framework.openai import OpenAIChatClient
    from openai import AsyncOpenAI
    AI_AVAILABLE = True
except ImportError:
    AI_AVAILABLE = False
    print("⚠️  AI Agent Framework not available. Only direct generation will be supported.")

# PowerPoint imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE


class UnifiedPowerPointGenerator:
    """Unified PowerPoint generator with both AI and direct generation capabilities."""
    
    def __init__(self, github_token: str = None):
        """Initialize the unified generator.
        
        Args:
            github_token: GitHub personal access token for AI model access (optional)
        """
        self.github_token = github_token
        self.ai_agent = None
        
        # Conceptual building blocks (no Azure product names)
        self.conceptual_blocks = {
            "web_application": {
                "name": "User Experience & Application Layer",
                "concepts": ["Web Portals", "Mobile Applications", "User Interfaces", "Frontend Services"]
            },
            "ai_analytics": {
                "name": "Data & Intelligence Layer", 
                "concepts": ["AI/ML Models", "Analytics Engine", "Data Processing", "Intelligent Services"]
            },
            "data_platform": {
                "name": "Data Platform Layer",
                "concepts": ["Databases", "Data Storage", "Data Lakes", "Data Pipelines"]
            },
            "integration": {
                "name": "Integration & API Layer",
                "concepts": ["API Gateway", "Message Queues", "Workflow Engine", "Event Processing"]
            }
        }
        
        # Azure service recommendations by category
        self.service_recommendations = {
            "web_application": {
                "name": "User Experience & Application Layer",
                "services": ["Azure Web Apps", "Azure Container Apps", "Azure Kubernetes Service", "Azure Front Door"]
            },
            "ai_analytics": {
                "name": "Data & Intelligence Layer", 
                "services": ["Azure OpenAI", "Microsoft Fabric", "Azure Databricks", "Azure AI Services"]
            },
            "data_platform": {
                "name": "Data Platform Layer",
                "services": ["Azure SQL Database", "Azure Cosmos DB", "Azure Storage", "Azure Data Factory"]
            },
            "integration": {
                "name": "Integration & API Layer",
                "services": ["Azure API Management", "Azure Service Bus", "Azure Logic Apps", "Azure Event Grid"]
            }
        }
        
        # Cross-cutting services (always included)
        self.conceptual_cross_cutting = [
            "Security & Compliance", "Network Security", "Monitoring & Defense",
            "Data Encryption", "System Monitoring", "Backup & Recovery", "Identity Management"
        ]
        
        self.azure_cross_cutting_services = [
            "Azure Policy and Compliance", "Azure Firewall and DDoS", "Microsoft Sentinel and Defender",
            "Encryption", "Azure Monitor", "Azure Backup and BCDR", "Microsoft Entra ID"
        ]
        
        # Azure service icons mapping (using Unicode symbols as placeholders for actual icons)
        self.azure_icons = {
            # Web Services
            "Azure Web Apps": "🌐",
            "Azure Container Apps": "📦",
            "Azure Kubernetes Service": "⚙️",
            "Azure Front Door": "🚪",
            "Azure App Service": "🌐",
            "Azure Application Gateway": "🚪",
            "Azure Load Balancer": "⚖️",
            
            # AI/ML Services
            "Azure OpenAI": "🤖",
            "Microsoft Fabric": "🧩",
            "Azure Databricks": "📊",
            "Azure AI Services": "🧠",
            "Azure Synapse Analytics": "📈",
            "Azure ML": "🔬",
            "Power BI": "📊",
            "Azure AI Search": "🔍",
            
            # Data Services
            "Azure SQL Database": "🗄️",
            "Azure Cosmos DB": "🌌",
            "Azure Storage": "💾",
            "Azure Data Factory": "🏭",
            "Azure Data Lake": "🏞️",
            "Azure Synapse": "📈",
            "Azure Purview": "🔍",
            
            # Integration Services
            "Azure API Management": "🔌",
            "Azure Service Bus": "🚌",
            "Azure Logic Apps": "⚡",
            "Azure Event Grid": "📋",
            "Azure Event Hub": "📡",
            "Function Apps": "⚡",
            "Power Automate": "🔄",
            
            # Security Services
            "Microsoft Entra ID": "🔐",
            "Azure Key Vault": "🔑",
            "Microsoft Sentinel": "🛡️",
            "Azure Firewall": "🔥",
            "Microsoft Defender": "🛡️",
            "Azure Policy": "📋",
            "Encryption": "🔒",
            
            # Infrastructure Services
            "Azure Virtual Networks": "🌐",
            "Azure Virtual Machines": "🖥️",
            "Azure Monitor": "📊",
            "Azure DevOps": "🔧",
            "Azure Backup": "💾",
            "Azure Policy and Compliance": "📋",
            "Azure Firewall and DDoS": "🔥",
            "Microsoft Sentinel and Defender": "🛡️",
            "Azure Backup and BCDR": "💾"
        }
        
        # Building block colors
        self.colors = {
            "web_application": RGBColor(0, 120, 212),    # Azure Blue
            "ai_analytics": RGBColor(138, 43, 226),      # Purple  
            "data_platform": RGBColor(0, 188, 140),      # Teal
            "integration": RGBColor(255, 140, 0),        # Orange
            "security": RGBColor(232, 17, 35),           # Red
            "infrastructure": RGBColor(16, 110, 190)     # Dark Blue
        }
        
        # Extended Azure building blocks for comprehensive presentations
        self.azure_building_blocks = {
            "compute": {
                "virtual_machines": "Azure Virtual Machines (VMs) provide on-demand, scalable computing resources",
                "app_service": "Azure App Service hosts web apps, REST APIs, and mobile backends",
                "azure_functions": "Azure Functions enables serverless compute for event-driven applications",
                "container_instances": "Azure Container Instances run containers without managing servers",
                "kubernetes_service": "Azure Kubernetes Service (AKS) manages containerized applications"
            },
            "storage": {
                "blob_storage": "Azure Blob Storage stores massive amounts of unstructured object data",
                "disk_storage": "Azure Disk Storage provides high-performance, durable block storage",
                "files": "Azure Files offers fully managed file shares in the cloud",
                "queue_storage": "Azure Queue Storage provides messaging between application components",
                "table_storage": "Azure Table Storage stores structured NoSQL data"
            },
            "networking": {
                "virtual_network": "Azure Virtual Network enables secure communication between Azure resources",
                "load_balancer": "Azure Load Balancer distributes incoming traffic across healthy VMs",
                "application_gateway": "Azure Application Gateway provides application-level routing and load balancing",
                "vpn_gateway": "Azure VPN Gateway connects on-premises networks to Azure",
                "express_route": "Azure ExpressRoute creates private connections to Azure datacenters"
            },
            "database": {
                "sql_database": "Azure SQL Database provides managed relational database service",
                "cosmos_db": "Azure Cosmos DB offers globally distributed, multi-model database service",
                "postgresql": "Azure Database for PostgreSQL provides managed PostgreSQL service",
                "mysql": "Azure Database for MySQL offers managed MySQL service",
                "redis_cache": "Azure Cache for Redis provides in-memory data caching"
            },
            "ai_ml": {
                "cognitive_services": "Azure Cognitive Services provides AI capabilities via REST APIs",
                "machine_learning": "Azure Machine Learning enables building and deploying ML models",
                "ai_search": "Azure AI Search provides full-text search capabilities",
                "openai_service": "Azure OpenAI Service offers OpenAI models with enterprise security",
                "bot_service": "Azure Bot Service builds conversational AI experiences"
            },
            "security": {
                "key_vault": "Azure Key Vault securely stores and manages secrets, keys, and certificates",
                "active_directory": "Azure Active Directory provides identity and access management",
                "security_center": "Azure Security Center provides unified security management",
                "sentinel": "Azure Sentinel offers cloud-native SIEM and SOAR capabilities",
                "firewall": "Azure Firewall provides network security filtering"
            },
            "monitoring": {
                "monitor": "Azure Monitor provides comprehensive monitoring for applications and infrastructure",
                "log_analytics": "Azure Log Analytics collects and analyzes log data",
                "application_insights": "Azure Application Insights monitors live applications",
                "service_health": "Azure Service Health provides insights into Azure service issues",
                "advisor": "Azure Advisor provides personalized recommendations"
            }
        }

    # =============================================================================
    # DIRECT GENERATION METHODS
    # =============================================================================
    
    def analyze_requirements(self, requirements: str) -> List[str]:
        """Analyze requirements and return needed categories."""
        req_lower = requirements.lower()
        categories = []
        
        # Check for layer mentions
        if any(term in req_lower for term in ["user experience", "ux", "frontend", "application layer", "app layer", "web", "portal"]):
            categories.append("web_application")
        
        if any(term in req_lower for term in ["data and intelligence", "data intelligence", "analytics", "ai", "machine learning"]):
            categories.append("ai_analytics")
            if "data" in req_lower:
                categories.append("data_platform")
        
        if any(term in req_lower for term in ["integration layer", "integration", "api"]):
            categories.append("integration")
        
        # Remove duplicates while preserving order
        return list(dict.fromkeys(categories))
    
    def get_azure_icon_symbol(self, service_name: str) -> str:
        """Get the icon symbol for an Azure service."""
        return self.azure_icons.get(service_name, "🔧")  # Default to gear icon
    
    def load_azure_icon_image(self, slide, x_pos, y_pos, service_name, icon_folder="azure_icons"):
        """Load actual Azure service icon if available (SVG or PNG)."""
        # Check if icon folder exists
        if not os.path.exists(icon_folder):
            return None
        
        # Common Azure service icon filename patterns
        possible_filenames = [
            f"{service_name.lower().replace(' ', '-')}.svg",
            f"{service_name.lower().replace(' ', '_')}.svg",
            f"{service_name.lower().replace(' ', '-')}.png",
            f"{service_name.lower().replace(' ', '_')}.png",
            f"{service_name.lower()}.svg",
            f"{service_name.lower()}.png"
        ]
        
        for filename in possible_filenames:
            icon_path = os.path.join(icon_folder, filename)
            if os.path.exists(icon_path):
                try:
                    # Add image to slide
                    icon_shape = slide.shapes.add_picture(
                        icon_path,
                        Inches(x_pos), Inches(y_pos),
                        Inches(0.3), Inches(0.3)  # Small icon size
                    )
                    print(f"Loaded Azure icon: {filename}")
                    return icon_shape
                except Exception as e:
                    print(f"Error loading icon {filename}: {e}")
                    continue
        
        return None
    
    def create_conceptual_slide(self, prs, categories, requirements):
        """Create the first slide with conceptual building blocks (no Azure product names)."""
        # Create blank slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Conceptual Architecture Building Blocks"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        
        # Calculate layout for conceptual blocks
        num_categories = len(categories)
        if num_categories <= 3:
            blocks_per_row = num_categories
            block_width = 2.8
            block_height = 2.2
        else:
            blocks_per_row = 3
            block_width = 2.5
            block_height = 2.0
        
        start_x = 0.5
        start_y = 1.3
        spacing_x = 0.3
        spacing_y = 0.4
        
        # Create conceptual building blocks
        for i, category in enumerate(categories):
            if category in self.conceptual_blocks:
                row = i // blocks_per_row
                col = i % blocks_per_row
                
                x_pos = start_x + col * (block_width + spacing_x)
                y_pos = start_y + row * (block_height + spacing_y)
                
                # Create block
                block_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(x_pos), Inches(y_pos),
                    Inches(block_width), Inches(block_height)
                )
                
                # Set color (lighter/pastel versions)
                fill = block_shape.fill
                fill.solid()
                fill.fore_color.rgb = self.colors[category]
                
                # Add text
                text_frame = block_shape.text_frame
                text_frame.clear()
                text_frame.margin_top = Inches(0.1)
                text_frame.margin_bottom = Inches(0.1)
                text_frame.margin_left = Inches(0.1)
                text_frame.margin_right = Inches(0.1)
                text_frame.word_wrap = True
                
                # Category title
                p_title = text_frame.paragraphs[0]
                p_title.text = self.conceptual_blocks[category]["name"]
                p_title.font.bold = True
                p_title.font.size = Pt(12)
                p_title.font.color.rgb = RGBColor(255, 255, 255)
                p_title.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                
                # Add conceptual services
                concepts = self.conceptual_blocks[category]["concepts"]
                for concept in concepts[:4]:  # Show top 4 concepts
                    p_concept = text_frame.add_paragraph()
                    p_concept.text = f"• {concept}"
                    p_concept.font.size = Pt(9)
                    p_concept.font.color.rgb = RGBColor(255, 255, 255)
                    p_concept.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    p_concept.space_after = Pt(3)
        
        # Add conceptual cross-cutting services as black boxes
        security_y_pos = start_y + ((num_categories - 1) // blocks_per_row + 1) * (block_height + spacing_y) + 0.3
        
        box_width = 1.25
        box_height = 0.6
        start_x_cross = 0.5
        spacing_x_cross = 0.05
        
        for i, service in enumerate(self.conceptual_cross_cutting):
            x_pos = start_x_cross + i * (box_width + spacing_x_cross)
            
            # Create rectangular box (not rounded)
            service_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x_pos), Inches(security_y_pos),
                Inches(box_width), Inches(box_height)
            )
            
            # Set black background
            fill = service_box.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)  # Pure black background
            
            # Add service text to the box
            service_text_frame = service_box.text_frame
            service_text_frame.clear()
            service_text_frame.margin_top = Inches(0.05)
            service_text_frame.margin_bottom = Inches(0.05)
            service_text_frame.margin_left = Inches(0.05)
            service_text_frame.margin_right = Inches(0.05)
            service_text_frame.word_wrap = True
            
            # Service text
            p_service = service_text_frame.paragraphs[0]
            p_service.text = service
            p_service.font.size = Pt(8)
            p_service.font.color.rgb = RGBColor(255, 255, 255)  # White text on black background
            p_service.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        
        # Add header label above the cross-cutting services
        header_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(security_y_pos - 0.3), Inches(9), Inches(0.25)
        )
        header_frame = header_box.text_frame
        header_frame.text = "Cross-cutting Security and Infrastructure Services"
        header_frame.paragraphs[0].font.bold = True
        header_frame.paragraphs[0].font.size = Pt(11)
        header_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        
        return security_y_pos, box_height

    def create_azure_specific_slide(self, prs, categories, requirements):
        """Create the second slide with Azure-specific services."""
        # Create blank slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Azure Solution Architecture Building Blocks"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        
        # Calculate layout
        num_categories = len(categories)
        if num_categories <= 3:
            blocks_per_row = num_categories
            block_width = 2.8
            block_height = 2.2
        else:
            blocks_per_row = 3
            block_width = 2.5
            block_height = 2.0
        
        start_x = 0.5
        start_y = 1.3
        spacing_x = 0.3
        spacing_y = 0.4
        
        # Create building blocks
        for i, category in enumerate(categories):
            if category in self.service_recommendations:
                row = i // blocks_per_row
                col = i % blocks_per_row
                
                x_pos = start_x + col * (block_width + spacing_x)
                y_pos = start_y + row * (block_height + spacing_y)
                
                # Create block
                block_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(x_pos), Inches(y_pos),
                    Inches(block_width), Inches(block_height)
                )
                
                # Set color
                fill = block_shape.fill
                fill.solid()
                fill.fore_color.rgb = self.colors[category]
                
                # Add text
                text_frame = block_shape.text_frame
                text_frame.clear()
                text_frame.margin_top = Inches(0.1)
                text_frame.margin_bottom = Inches(0.1)
                text_frame.margin_left = Inches(0.1)
                text_frame.margin_right = Inches(0.1)
                text_frame.word_wrap = True
                
                # Category title
                p_title = text_frame.paragraphs[0]
                p_title.text = self.service_recommendations[category]["name"]
                p_title.font.bold = True
                p_title.font.size = Pt(12)
                p_title.font.color.rgb = RGBColor(255, 255, 255)
                p_title.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                
                # Add Azure services with icons
                services = self.service_recommendations[category]["services"]
                for j, service in enumerate(services[:4]):  # Show top 4 services
                    p_service = text_frame.add_paragraph()
                    
                    # Try to load actual Azure icon first, fallback to Unicode symbol
                    icon_loaded = self.load_azure_icon_image(
                        slide, (block_shape.left / 914400) - 0.35, 
                        (block_shape.top / 914400) + (len(text_frame.paragraphs) * 0.15), service
                    )
                    
                    if icon_loaded:
                        # If we loaded an actual icon, just use the service name
                        p_service.text = service
                    else:
                        # Fallback to Unicode symbol + service name
                        icon_symbol = self.get_azure_icon_symbol(service)
                        p_service.text = f"{icon_symbol} {service}"
                    
                    p_service.font.size = Pt(9)
                    p_service.font.color.rgb = RGBColor(255, 255, 255)
                    p_service.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    p_service.space_after = Pt(3)
        
        # Add Azure cross-cutting services as black boxes
        security_y_pos = start_y + ((num_categories - 1) // blocks_per_row + 1) * (block_height + spacing_y) + 0.3
        
        box_width = 1.25
        box_height = 0.6
        start_x_cross = 0.5
        spacing_x_cross = 0.05
        
        for i, service in enumerate(self.azure_cross_cutting_services):
            x_pos = start_x_cross + i * (box_width + spacing_x_cross)
            
            # Create rectangular box (not rounded)
            service_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x_pos), Inches(security_y_pos),
                Inches(box_width), Inches(box_height)
            )
            
            # Set black background
            fill = service_box.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)  # Pure black background
            
            # Add service text to the box
            service_text_frame = service_box.text_frame
            service_text_frame.clear()
            service_text_frame.margin_top = Inches(0.05)
            service_text_frame.margin_bottom = Inches(0.05)
            service_text_frame.margin_left = Inches(0.05)
            service_text_frame.margin_right = Inches(0.05)
            service_text_frame.word_wrap = True
            
            # Service text with icon
            p_service = service_text_frame.paragraphs[0]
            
            # Try to load actual Azure icon first, fallback to Unicode symbol
            icon_loaded = self.load_azure_icon_image(
                slide, (service_box.left / 914400), (service_box.top / 914400) - 0.2, service
            )
            
            if icon_loaded:
                # If we loaded an actual icon, just use the service name
                p_service.text = service
            else:
                # Fallback to Unicode symbol + service name
                icon_symbol = self.get_azure_icon_symbol(service)
                p_service.text = f"{icon_symbol}\n{service}"
            
            p_service.font.size = Pt(7)
            p_service.font.color.rgb = RGBColor(255, 255, 255)  # White text on black background
            p_service.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        
        # Add header label above the cross-cutting services
        header_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(security_y_pos - 0.3), Inches(9), Inches(0.25)
        )
        header_frame = header_box.text_frame
        header_frame.text = "Cross-cutting Security and Infrastructure Services"
        header_frame.paragraphs[0].font.bold = True
        header_frame.paragraphs[0].font.size = Pt(11)
        header_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        
        return security_y_pos, box_height

    def create_building_block_presentation(self, requirements: str, filename: str = None, slide_type: str = "both") -> str:
        """Create building block presentation with conceptual and/or Azure-specific slides.
        
        Args:
            requirements: The requirements for the architecture
            filename: Name of the output file
            slide_type: "conceptual", "azure", or "both" (default)
        """
        try:
            # Analyze requirements
            categories = self.analyze_requirements(requirements)
            
            print(f"Identified categories: {[cat.replace('_', ' ').title() for cat in categories]}")
            
            # Create presentation
            prs = Presentation()
            # Remove default slide
            if len(prs.slides) > 0:
                slide_to_remove = prs.slides[0]
                rId = prs.slides.slides._element.index(slide_to_remove._element)
                prs.part.drop_rel(prs.slides._sld_id_lst[rId].rId)
                del prs.slides._sld_id_lst[rId]
            
            security_y_pos = None
            box_height = None
            
            # Create slides based on type
            if slide_type in ["conceptual", "both"]:
                print("Creating conceptual building blocks slide...")
                security_y_pos1, box_height1 = self.create_conceptual_slide(prs, categories, requirements)
                security_y_pos = security_y_pos1
                box_height = box_height1
            
            if slide_type in ["azure", "both"]:
                print("Creating Azure-specific building blocks slide...")
                security_y_pos2, box_height2 = self.create_azure_specific_slide(prs, categories, requirements)
                if security_y_pos is None:
                    security_y_pos = security_y_pos2
                    box_height = box_height2
            
            # Add requirements at bottom of all slides
            for slide_num, slide in enumerate(prs.slides, 1):
                req_y_pos = security_y_pos + box_height + 0.2
                req_box = slide.shapes.add_textbox(Inches(0.5), Inches(req_y_pos), Inches(9), Inches(0.5))
                req_frame = req_box.text_frame
                req_frame.text = f"Requirements: {requirements}"
                req_frame.paragraphs[0].font.size = Pt(10)
                req_frame.paragraphs[0].font.italic = True
                req_frame.word_wrap = True
            
            # Save presentation
            if not filename:
                filename = "Building_Blocks_Architecture.pptx"
            elif not filename.endswith('.pptx'):
                filename += '.pptx'
            
            filepath = os.path.join(os.getcwd(), filename)
            prs.save(filepath)
            
            slide_count = len(prs.slides)
            slide_description = []
            if slide_type in ["conceptual", "both"]:
                slide_description.append("Conceptual Architecture Building Blocks")
            if slide_type in ["azure", "both"]:
                slide_description.append("Azure-Specific Architecture Building Blocks")
            
            return f"✅ Successfully created {slide_count} slide(s) with {len(categories)} building blocks\n   " + "\n   ".join([f"Slide {i+1}: {desc}" for i, desc in enumerate(slide_description)]) + f"\nSaved as: {filename}"
            
        except Exception as e:
            return f"❌ Error creating presentation: {str(e)}"

    # =============================================================================
    # AI-POWERED GENERATION METHODS (when available)
    # =============================================================================
    
    async def initialize_ai_agent(self, model_id: str = "openai/gpt-4o-mini"):
        """Initialize the AI agent with tools (if available)."""
        if not AI_AVAILABLE or not self.github_token:
            return False
        
        try:
            openai_client = AsyncOpenAI(
                base_url="https://models.github.ai/inference",
                api_key=self.github_token,
            )
            
            chat_client = OpenAIChatClient(
                async_client=openai_client,
                model_id=model_id
            )
            
            self.ai_agent = ChatAgent(
                chat_client=chat_client,
                name="UnifiedPowerPointAgent",
                instructions="""You are an expert PowerPoint creation assistant specializing in Azure architecture.

Your role is to:
1. Analyze user requests and identify the appropriate solution categories needed
2. Recommend suitable Azure services based on requirements
3. Create professional PowerPoint presentations with building block architecture
4. Generate both conceptual and Azure-specific slides when requested

Focus on:
- Clear, professional slide layouts
- Logical architecture organization  
- Best practices for Azure services
- Building block format with colored sections
- Cross-cutting security and infrastructure services

Use the available tools to create comprehensive presentations.""",
                tools=[
                    self.ai_get_azure_services,
                    self.ai_create_powerpoint_slide,
                    self.ai_create_architecture_diagram,
                    self.ai_save_presentation
                ]
            )
            return True
        except Exception as e:
            print(f"Failed to initialize AI agent: {e}")
            return False
    
    def ai_get_azure_services(
        self, 
        category: Annotated[str, "Category of Azure services (compute, storage, networking, database, ai_ml, security, monitoring)"]
    ) -> str:
        """Get information about Azure services in a specific category."""
        if category.lower() in self.azure_building_blocks:
            services = self.azure_building_blocks[category.lower()]
            result = f"Azure {category.title()} Services:\n\n"
            for service, description in services.items():
                result += f"• {service.replace('_', ' ').title()}: {description}\n"
            return result
        else:
            available_categories = ", ".join(self.azure_building_blocks.keys())
            return f"Category '{category}' not found. Available categories: {available_categories}"
    
    def ai_create_powerpoint_slide(
        self,
        title: Annotated[str, "Title of the slide"],
        content: Annotated[str, "Main content for the slide"],
        slide_type: Annotated[str, "Type of slide: title, content, or building_block"] = "content"
    ) -> str:
        """Create a PowerPoint slide with the specified content."""
        try:
            # Initialize presentation if it doesn't exist
            if not hasattr(self, 'ai_presentation'):
                self.ai_presentation = Presentation()
                # Remove the default slide
                if len(self.ai_presentation.slides) > 0:
                    slide_to_remove = self.ai_presentation.slides[0]
                    rId = self.ai_presentation.slides.slides._element.index(slide_to_remove._element)
                    self.ai_presentation.part.drop_rel(self.ai_presentation.slides._sld_id_lst[rId].rId)
                    del self.ai_presentation.slides._sld_id_lst[rId]
            
            # Handle different slide types
            if slide_type.lower() == "building_block":
                # Use the direct generation method for building blocks
                categories = self.analyze_requirements(content)
                prs = self.ai_presentation
                security_y_pos, box_height = self.create_azure_specific_slide(prs, categories, content)
                return f"Successfully created building block slide: '{title}'"
            
            # Standard slide creation
            if slide_type.lower() == "title":
                slide_layout = self.ai_presentation.slide_layouts[0]  # Title slide
                slide = self.ai_presentation.slides.add_slide(slide_layout)
                slide.shapes.title.text = title
                if slide.shapes.placeholders[1]:  # Subtitle
                    slide.shapes.placeholders[1].text = content
            else:
                slide_layout = self.ai_presentation.slide_layouts[1]  # Title and content
                slide = self.ai_presentation.slides.add_slide(slide_layout)
                slide.shapes.title.text = title
                
                # Add content to the body
                content_placeholder = slide.shapes.placeholders[1]
                content_placeholder.text = content
                
                # Format the text
                for paragraph in content_placeholder.text_frame.paragraphs:
                    paragraph.font.size = Pt(14)
                    paragraph.font.name = 'Segoe UI'
            
            slide_number = len(self.ai_presentation.slides)
            return f"Successfully created slide {slide_number}: '{title}'"
            
        except Exception as e:
            return f"Error creating slide: {str(e)}"
    
    def ai_create_architecture_diagram(
        self,
        title: Annotated[str, "Title of the architecture slide"],
        components: Annotated[List[str], "List of Azure components to include in the diagram"],
        description: Annotated[str, "Description of the architecture"]
    ) -> str:
        """Create an architecture diagram slide with Azure components."""
        try:
            # Initialize presentation if it doesn't exist
            if not hasattr(self, 'ai_presentation'):
                self.ai_presentation = Presentation()
                if len(self.ai_presentation.slides) > 0:
                    slide_to_remove = self.ai_presentation.slides[0]
                    rId = self.ai_presentation.slides.slides._element.index(slide_to_remove._element)
                    self.ai_presentation.part.drop_rel(self.ai_presentation.slides._sld_id_lst[rId].rId)
                    del self.ai_presentation.slides._sld_id_lst[rId]
            
            # Create blank slide for architecture diagram
            slide_layout = self.ai_presentation.slide_layouts[6]  # Blank layout
            slide = self.ai_presentation.slides.add_slide(slide_layout)
            
            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = Pt(24)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            
            # Create simple component boxes
            start_x = 1.5
            start_y = 2.0
            box_width = 2.0
            box_height = 1.0
            spacing = 0.5
            
            colors = [
                RGBColor(0, 120, 212),    # Azure blue
                RGBColor(0, 188, 140),    # Green
                RGBColor(255, 140, 0),    # Orange
                RGBColor(232, 17, 35),    # Red
                RGBColor(136, 23, 152),   # Purple
                RGBColor(16, 110, 190),   # Dark blue
            ]
            
            for i, component in enumerate(components[:6]):  # Limit to 6 components
                x_pos = start_x + (i % 3) * (box_width + spacing)
                y_pos = start_y + (i // 3) * (box_height + spacing)
                
                # Create rectangle shape
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(x_pos), Inches(y_pos),
                    Inches(box_width), Inches(box_height)
                )
                
                # Set fill color
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = colors[i % len(colors)]
                
                # Add text
                text_frame = shape.text_frame
                text_frame.text = component.replace('_', ' ').title()
                text_frame.text_frame_format.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                paragraph = text_frame.paragraphs[0]
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.font.bold = True
                paragraph.font.size = Pt(12)
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            
            # Add description box
            desc_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(5.5), Inches(9), Inches(1.5)
            )
            desc_frame = desc_box.text_frame
            desc_frame.text = description
            desc_frame.paragraphs[0].font.size = Pt(14)
            desc_frame.word_wrap = True
            
            slide_number = len(self.ai_presentation.slides)
            return f"Successfully created architecture diagram slide {slide_number}: '{title}' with {len(components)} components"
            
        except Exception as e:
            return f"Error creating architecture diagram: {str(e)}"
    
    def ai_save_presentation(
        self,
        filename: Annotated[str, "Name of the PowerPoint file to save (without .pptx extension)"]
    ) -> str:
        """Save the AI-generated PowerPoint presentation to a file."""
        try:
            if not hasattr(self, 'ai_presentation'):
                return "No presentation to save. Create slides first."
            
            # Ensure filename ends with .pptx
            if not filename.endswith('.pptx'):
                filename += '.pptx'
            
            # Save to the current directory
            filepath = os.path.join(os.getcwd(), filename)
            self.ai_presentation.save(filepath)
            
            slide_count = len(self.ai_presentation.slides)
            return f"Successfully saved presentation '{filename}' with {slide_count} slides to {filepath}"
            
        except Exception as e:
            return f"Error saving presentation: {str(e)}"
    
    async def create_ai_presentation(self, user_request: str) -> str:
        """Process user request using AI to create a PowerPoint presentation."""
        if not AI_AVAILABLE:
            return "❌ AI functionality not available. Please install agent-framework-azure-ai."
        
        if not self.github_token:
            return "❌ GitHub token required for AI functionality."
        
        # Initialize AI agent if not done
        if not self.ai_agent:
            success = await self.initialize_ai_agent()
            if not success:
                return "❌ Failed to initialize AI agent."
        
        thread = self.ai_agent.get_new_thread()
        
        print("🤖 AI PowerPoint Agent: ", end="", flush=True)
        response_text = ""
        
        async for chunk in self.ai_agent.run_stream(user_request, thread=thread):
            if chunk.text:
                print(chunk.text, end="", flush=True)
                response_text += chunk.text
        
        print("\n")
        return response_text

    # =============================================================================
    # UNIFIED INTERFACE METHODS
    # =============================================================================
    
    def generate(self, requirements: str, method: str = "direct", filename: str = None, 
                 slide_type: str = "both") -> str:
        """Main generation method that routes to appropriate implementation.
        
        Args:
            requirements: The requirements for the architecture
            method: "direct", "ai", or "auto" (chooses best available)
            filename: Name of the output file
            slide_type: For direct method: "conceptual", "azure", or "both"
        """
        if method == "auto":
            method = "ai" if AI_AVAILABLE and self.github_token else "direct"
        
        if method == "ai":
            if not AI_AVAILABLE:
                print("⚠️  AI method not available, falling back to direct generation")
                method = "direct"
            elif not self.github_token:
                print("⚠️  GitHub token required for AI method, falling back to direct generation")
                method = "direct"
        
        if method == "direct":
            return self.create_building_block_presentation(requirements, filename, slide_type)
        elif method == "ai":
            # For AI method, we run async
            return asyncio.run(self.create_ai_presentation(requirements))
        else:
            return f"❌ Unknown method: {method}. Use 'direct', 'ai', or 'auto'."


def main():
    """Main function with interactive menu."""
    # Get GitHub token from environment
    github_token = os.getenv('GITHUB_TOKEN')
    
    # Initialize generator
    generator = UnifiedPowerPointGenerator(github_token)
    
    print("🚀 Unified PowerPoint Generator")
    print("=" * 60)
    print("Combines all PowerPoint generation capabilities in one tool!")
    print()
    
    # Show capabilities
    print("📋 Available Generation Methods:")
    print("1. Direct Generation - Fast, reliable building block slides")
    print("2. AI-Powered Generation - Intelligent content creation")
    if not AI_AVAILABLE:
        print("   ⚠️  AI features require: pip install agent-framework-azure-ai --pre")
    if not github_token:
        print("   ⚠️  AI features require GitHub token in GITHUB_TOKEN environment variable")
    print()
    
    print("🏗️  Service Categories Supported:")
    for category, data in generator.service_recommendations.items():
        services = ", ".join(data["services"][:3])  # Show first 3 services
        print(f"   • {data['name']}: {services}...")
    print()
    
    # Interactive loop
    while True:
        try:
            print("=" * 60)
            print("Choose generation method:")
            print("1. Direct Generation (Building Blocks)")
            print("2. AI-Powered Generation (Full Presentations)")
            print("3. Quick Test (Your example requirements)")
            print("4. Exit")
            print()
            
            choice = input("Select option (1-4): ").strip()
            
            if choice == '4':
                print("👋 Goodbye!")
                break
            
            if choice == '3':
                # Quick test with your example
                test_req = "1. User experience layer, 2. Application Layer, 3. Data and intelligence layer, 4. Integration Layer"
                print(f"\n🧪 Testing with: {test_req}")
                result = generator.generate(test_req, method="direct", filename="Test_Multi_Layer_Architecture")
                print(result)
                continue
            
            # Get requirements
            print("\n💡 Example requirements:")
            examples = [
                "AI-powered customer service with analytics and web interface",
                "E-commerce platform with payment processing and inventory management",
                "Healthcare management system with patient portal and data analytics",
                "Financial services platform with fraud detection and compliance"
            ]
            for i, example in enumerate(examples, 1):
                print(f"   {i}. {example}")
            print()
            
            requirements = input("Enter your requirements: ").strip()
            if not requirements:
                print("❌ Requirements cannot be empty.")
                continue
            
            filename = input("Output filename (optional, press Enter for auto): ").strip()
            if not filename:
                filename = None
            
            print(f"\n🔄 Generating presentation...")
            
            if choice == '1':
                # Direct generation
                print("Choose slide type:")
                print("1. Both (Conceptual + Azure-specific)")
                print("2. Conceptual only")
                print("3. Azure-specific only")
                
                slide_choice = input("Select (1-3, default=1): ").strip() or "1"
                slide_types = {"1": "both", "2": "conceptual", "3": "azure"}
                slide_type = slide_types.get(slide_choice, "both")
                
                result = generator.generate(requirements, method="direct", 
                                           filename=filename, slide_type=slide_type)
                print(result)
                
            elif choice == '2':
                # AI generation
                if not AI_AVAILABLE or not github_token:
                    print("❌ AI generation not available. Using direct generation instead.")
                    result = generator.generate(requirements, method="direct", filename=filename)
                    print(result)
                else:
                    result = generator.generate(requirements, method="ai", filename=filename)
                    print(result)
            else:
                print("❌ Invalid choice. Please select 1-4.")
        
        except KeyboardInterrupt:
            print("\n👋 Goodbye!")
            break
        except Exception as e:
            print(f"❌ Error: {str(e)}")


if __name__ == "__main__":
    main()